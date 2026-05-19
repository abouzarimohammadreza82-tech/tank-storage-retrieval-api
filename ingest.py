import os

# =========================
# تنظیم API Key برای Pinecone
# =========================

os.environ["PINECONE_API_KEY"] = "pcsk_6XTJiF_Azamvb6rnwqtZQYaabxDvuaYj794YfSEG39LSvbb6ZkL5tcEf9itQNFp7t81QwM"

from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

from pinecone import Pinecone
from langchain_pinecone import PineconeVectorStore


# =========================
# تنظیمات
# =========================

DATA_DIR = "data"

INDEX_NAME = "tank-storage-db"


# =========================
# استخراج متن از پاورپوینت
# =========================

def extract_text_from_pptx(filepath):

    prs = Presentation(filepath)

    slides_text = []

    for i, slide in enumerate(prs.slides):

        slide_content = []

        for shape in slide.shapes:

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:

                text = para.text.strip()

                if text:
                    slide_content.append(text)

        if slide_content:

            full_text = (
                f"[اسلاید {i+1}]\n"
                + "\n".join(slide_content)
            )

            slides_text.append(full_text)

    return slides_text


# =========================
# ساخت Vector Database
# =========================

def build_vector_db():

    all_texts = []

    # خواندن فایل‌های پاورپوینت

    for filename in os.listdir(DATA_DIR):

        if filename.endswith(".pptx"):

            filepath = os.path.join(DATA_DIR, filename)

            print(f"در حال پردازش: {filename}")

            texts = extract_text_from_pptx(filepath)

            all_texts.extend(texts)

    # اگر فایل پیدا نشد

    if not all_texts:

        print("هیچ فایل پاورپوینتی پیدا نشد!")

        return False

    print(f"تعداد اسلایدها: {len(all_texts)}")

    # =========================
    # Chunking
    # =========================

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )

    chunks = splitter.create_documents(all_texts)

    print(f"تعداد چانک‌ها: {len(chunks)}")

    # =========================
    # Embeddings
    # =========================

    print("در حال ساخت Embeddings...")

    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
    )

    # =========================
    # اتصال به Pinecone
    # =========================

    print("در حال اتصال به Pinecone...")

    pc = Pinecone()

    # تست اتصال به Index
    pc.Index(INDEX_NAME)

    # =========================
    # آپلود به Pinecone
    # =========================

    print("در حال آپلود داده‌ها به Pinecone...")

    PineconeVectorStore.from_documents(
        documents=chunks,
        embedding=embeddings,
        index_name=INDEX_NAME
    )

    print("Vector Database با موفقیت در Pinecone ذخیره شد!")

    return True


# =========================
# اجرای برنامه
# =========================

if __name__ == "__main__":

    build_vector_db()